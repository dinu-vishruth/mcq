# models/pdf_processor.py
import pdfplumber
from docx import Document
from pptx import Presentation

def extract_text_from_pdf(file_storage):
    try:
        file_storage.stream.seek(0)
    except Exception:
        pass
    text = ""
    with pdfplumber.open(file_storage) as pdf:
        for page in pdf.pages:
            t = page.extract_text()
            if t:
                text += t + "\n"
    return text

def extract_text_from_docx(path):
    doc = Document(path)
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

def extract_text_from_pptx(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
