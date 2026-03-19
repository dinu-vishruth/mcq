# models/pdf_processor.py
import pdfplumber
from docx import Document
from pptx import Presentation


def extract_text_from_pdf(file_storage):
    """
    Extract text from PDF including both paragraph text and table data.
    Reads ALL pages completely.
    """
    try:
        file_storage.stream.seek(0)
    except Exception:
        pass

    text_parts = []
    with pdfplumber.open(file_storage) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            # Extract regular text
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)

            # Extract tables (common in educational PDFs)
            tables = page.extract_tables()
            for table in tables:
                for row in table:
                    if row:
                        cleaned_row = [str(cell).strip() for cell in row if cell]
                        if cleaned_row:
                            text_parts.append(" | ".join(cleaned_row))

    return "\n".join(text_parts)


def extract_text_from_docx(path):
    """Extract text from DOCX including paragraphs and tables."""
    doc = Document(path)
    parts = []

    # Extract paragraphs
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)

    # Extract tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts)


def extract_text_from_pptx(path):
    """Extract text from PPTX slides including notes."""
    prs = Presentation(path)
    parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)

        # Also extract slide notes (often contain extra detail)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_text.append(notes)

        if slide_text:
            parts.append("\n".join(slide_text))

    return "\n\n".join(parts)
